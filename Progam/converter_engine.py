import os
import re
import csv

# ============================================================
# SAFE IMPORTS â€” Semua library opsional, tidak crash jika absen
# ============================================================
try:
    import docx
    import docx.text.paragraph
    import docx.table
except ImportError:
    docx = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None


# Ekstensi file gambar yang ditolak secara aman
IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp', 'tiff', 'tif', 'svg', 'ico', 'heic', 'heif'}

# Ekstensi yang dianggap source code
CODE_EXTENSIONS = {
    'py', 'js', 'ts', 'jsx', 'tsx', 'java', 'c', 'cpp', 'h', 'hpp', 'cs',
    'go', 'rs', 'rb', 'php', 'swift', 'kt', 'kts', 'r', 'sh', 'bat', 'ps1',
    'sql', 'json', 'yaml', 'yml', 'toml', 'xml', 'dart', 'scala', 'lua',
    'vim', 'makefile', 'dockerfile', 'gitignore', 'env', 'ini', 'cfg'
}


class MarkdownEngine:
    """
    CO-TO Core Engine v2 â€” Smart file-to-Markdown converter.
    Mendukung: TXT, MD, LOG, CSV, XLSX, DOCX (+ tabel), PDF,
               PPTX, HTML, Source Code. Dengan Universal Fallback & Image Guard.
    """

    def _clean_text(self, text: str) -> str:
        """Memotong spasi dan enter berlebih. Maksimal 2 baris kosong."""
        if not text:
            return ""
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = '\n'.join([line.rstrip() for line in text.split('\n')])
        return text.strip()

    def _clean_table_cell(self, value) -> str:
        """Normalisasi isi sel tabel agar aman untuk Markdown pipe table."""
        if value is None:
            return ""
        return str(value).replace('\n', ' ').strip()

    def _rows_to_markdown_table(self, rows) -> str:
        """Ubah list-of-lists menjadi Markdown table sederhana."""
        cleaned_rows = []
        max_cols = 0

        for row in rows:
            cleaned = [self._clean_table_cell(cell) for cell in (row or [])]
            if any(cleaned):
                cleaned_rows.append(cleaned)
                max_cols = max(max_cols, len(cleaned))

        if not cleaned_rows or max_cols == 0:
            return ""

        normalized_rows = [
            row + [""] * (max_cols - len(row))
            for row in cleaned_rows
        ]

        header = "| " + " | ".join(normalized_rows[0]) + " |"
        separator = "| " + " | ".join(["---"] * max_cols) + " |"
        body = [
            "| " + " | ".join(row) + " |"
            for row in normalized_rows[1:]
        ]

        return "\n".join([header, separator] + body)

    def _dataframe_to_markdown_table(self, df) -> str:
        """Ubah pandas DataFrame menjadi Markdown table tanpa tabulate."""
        df = df.fillna('')
        df = df.astype(str)

        rows = [df.columns.astype(str).tolist()]
        rows.extend(df.values.tolist())
        return self._rows_to_markdown_table(rows)

    # ----------------------------------------------------------
    # CONVERTER METHODS
    # ----------------------------------------------------------

    def _convert_txt(self, file_path: str) -> str:
        """Baca file teks biasa (txt, log, md, dsb.)."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        if self._looks_like_html(content):
            return self._convert_html(raw_html=content)
        return self._clean_text(content)

    def _convert_code(self, file_path: str, ext: str, filename: str) -> str:
        """Bungkus source code ke dalam fenced code block Markdown."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        # Jika content sudah mengandung triple-backtick, gunakan 4 backtick sebagai fence
        fence = '````' if '```' in content else '```'
        return f"# File: `{filename}`\n\n{fence}{ext}\n{content}\n{fence}"

    def _convert_docx(self, file_path: str) -> str:
        """
        Konversi DOCX ke Markdown.
        UPGRADE v2: Membaca elemen dokumen berurutan (paragraf + tabel Word asli).
        """
        if not docx:
            raise ImportError("Library 'python-docx' belum terinstall.")

        doc = docx.Document(file_path)
        md_lines = []

        for element in doc.element.body:

            # 1. PARAGRAF (Heading, List, atau teks biasa)
            if element.tag.endswith('p'):
                para = docx.text.paragraph.Paragraph(element, doc)
                text = para.text.strip()
                if not text:
                    continue

                style_name = para.style.name.lower()
                # Heading â€” bahasa Inggris & Indonesia
                if 'heading 1' in style_name or 'judul 1' in style_name:
                    md_lines.append(f"\n# {text}\n")
                elif 'heading 2' in style_name or 'judul 2' in style_name:
                    md_lines.append(f"\n## {text}\n")
                elif 'heading 3' in style_name or 'judul 3' in style_name:
                    md_lines.append(f"\n### {text}\n")
                elif 'heading 4' in style_name or 'judul 4' in style_name:
                    md_lines.append(f"\n#### {text}\n")
                # List item
                elif style_name.startswith('list') or text.startswith(('â€¢', '-', '*')):
                    clean_item = text.lstrip('â€¢-* ').strip()
                    md_lines.append(f"* {clean_item}")
                # Teks biasa
                else:
                    md_lines.append(text)

            # 2. TABEL WORD ASLI â†’ Markdown Table
            elif element.tag.endswith('tbl'):
                table = docx.table.Table(element, doc)
                md_lines.append("")  # Jarak sebelum tabel

                seen_rows = set()  # Hindari duplikasi sel merged
                for r_idx, row in enumerate(table.rows):
                    row_cells = []
                    for cell in row.cells:
                        cell_text = cell.text.strip().replace('\n', ' ')
                        row_cells.append(cell_text)

                    # Hapus duplikasi pada sel yang di-merge secara horizontal
                    # (python-docx mengembalikan teks sama untuk merged cells)
                    row_key = tuple(row_cells)
                    if row_key in seen_rows:
                        continue
                    seen_rows.add(row_key)

                    md_row = "| " + " | ".join(row_cells) + " |"
                    md_lines.append(md_row)

                    # Baris pertama = Header, tambahkan separator
                    if r_idx == 0:
                        divider = "| " + " | ".join(["---"] * len(row_cells)) + " |"
                        md_lines.append(divider)

                md_lines.append("")  # Jarak setelah tabel

        return self._clean_text("\n".join(md_lines))

    def _convert_pdf(self, file_path: str) -> str:
        """Ekstrak teks dan tabel dari setiap halaman PDF."""
        if not pdfplumber:
            raise ImportError("Library 'pdfplumber' belum terinstall.")

        md_lines = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                tables = page.extract_tables() or []
                table_blocks = []

                for table in tables:
                    md_table = self._rows_to_markdown_table(table)
                    if md_table:
                        table_blocks.append(md_table)

                if text or table_blocks:
                    md_lines.append(f"\n---\n*Halaman {i + 1}*\n")

                if text:
                    md_lines.append(text)

                for table_block in table_blocks:
                    md_lines.append("")
                    md_lines.append(table_block)

        return self._clean_text("\n".join(md_lines))

    def _extract_pptx_shape(self, shape, md_lines) -> bool:
        """Ekstrak tabel/teks PPTX, termasuk child shapes dalam group."""
        has_content = False

        if getattr(shape, 'has_table', False):
            table = shape.table
            md_lines.append("")
            for r_idx, row in enumerate(table.rows):
                row_cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                md_row = "| " + " | ".join(row_cells) + " |"
                md_lines.append(md_row)

                if r_idx == 0:
                    divider = "| " + " | ".join(["---"] * len(row.cells)) + " |"
                    md_lines.append(divider)
            md_lines.append("")
            has_content = True

        elif getattr(shape, 'has_text_frame', False):
            slide_texts = []
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    slide_texts.append(line)

            if slide_texts:
                md_lines.append(f"**{slide_texts[0]}**")
                for line in slide_texts[1:]:
                    md_lines.append(f"* {line}")
                has_content = True

        child_shapes = getattr(shape, 'shapes', None)
        if child_shapes:
            for child_shape in child_shapes:
                if self._extract_pptx_shape(child_shape, md_lines):
                    has_content = True

        return has_content

    def _convert_pptx(self, file_path: str) -> str:
        """
        Konversi PPTX ke Markdown.
        Mengiterasi setiap slide dan mengekstrak kotak teks biasa, tabel,
        serta isi grouped shapes secara berurutan.
        """
        if not PptxPresentation:
            raise ImportError("Library 'python-pptx' belum terinstall.")

        prs = PptxPresentation(file_path)
        md_lines = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            md_lines.append(f"\n---\n## Slide {slide_num}\n")
            has_content = False

            for shape in slide.shapes:
                if self._extract_pptx_shape(shape, md_lines):
                    has_content = True

            if not has_content:
                md_lines.append("*(Slide kosong, hanya berisi gambar, atau diagram grup)*")

        return self._clean_text("\n".join(md_lines))

    def _convert_tabular(self, file_path: str, ext: str) -> str:
        """
        Konversi CSV atau XLSX ke Markdown Table.

        v3:
        - CSV memakai csv module bawaan Python agar tidak bergantung pandas.
        - XLSX memakai openpyxl langsung dan mendukung multi-sheet.
        """
        MAX_ROWS = 500

        def clean_cell(value) -> str:
            if value is None:
                return ""
            text = str(value).replace("\n", " ").replace("\r", " ").strip()
            return text

        def escape_md_cell(value) -> str:
            # Escape pipe agar isi cell tidak merusak Markdown table.
            return clean_cell(value).replace("|", "\\|")

        def rows_to_markdown(rows, truncated_note: str = "") -> str:
            if not rows:
                return ""

            cleaned_rows = [
                [escape_md_cell(cell) for cell in row]
                for row in rows
            ]

            max_cols = max(len(row) for row in cleaned_rows)
            normalized = [
                row + [""] * (max_cols - len(row))
                for row in cleaned_rows
            ]

            headers = normalized[0]
            body_rows = normalized[1:]

            md_lines = []
            md_lines.append("| " + " | ".join(headers) + " |")
            md_lines.append("| " + " | ".join(["---"] * max_cols) + " |")

            for row in body_rows:
                md_lines.append("| " + " | ".join(row) + " |")

            if truncated_note:
                md_lines.append(truncated_note)

            return "\n".join(md_lines)

        if ext == "csv":
            last_error = None
            for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
                try:
                    with open(file_path, "r", encoding=encoding, errors="ignore", newline="") as f:
                        sample = f.read(4096)
                        f.seek(0)

                        try:
                            dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
                        except csv.Error:
                            dialect = csv.excel

                        reader = csv.reader(f, dialect)
                        rows = list(reader)

                    rows = [
                        row for row in rows
                        if any(clean_cell(cell) for cell in row)
                    ]

                    total_data_rows = max(0, len(rows) - 1)
                    truncated_note = ""
                    if total_data_rows > MAX_ROWS:
                        rows = rows[:MAX_ROWS + 1]
                        truncated_note = (
                            f"\n\n> Catatan: Data dipotong, hanya menampilkan "
                            f"{MAX_ROWS} dari {total_data_rows} baris untuk efisiensi token."
                        )

                    if not rows:
                        return ""

                    return self._clean_text(rows_to_markdown(rows, truncated_note))

                except Exception as e:
                    last_error = e

            raise ValueError(f"Gagal membaca CSV: {last_error}")

        if ext == "xlsx":
            if not load_workbook:
                raise ImportError("Library 'openpyxl' belum terinstall untuk membaca XLSX.")

            wb = load_workbook(file_path, read_only=True, data_only=True)
            md_sections = []

            try:
                for ws in wb.worksheets:
                    rows = []
                    for row in ws.iter_rows(values_only=True):
                        cleaned = [clean_cell(cell) for cell in row]

                        while cleaned and cleaned[-1] == "":
                            cleaned.pop()

                        if any(cleaned):
                            rows.append(cleaned)

                    if not rows:
                        continue

                    total_data_rows = max(0, len(rows) - 1)
                    truncated_note = ""
                    if total_data_rows > MAX_ROWS:
                        rows = rows[:MAX_ROWS + 1]
                        truncated_note = (
                            f"\n\n> Catatan: Sheet `{ws.title}` dipotong, hanya menampilkan "
                            f"{MAX_ROWS} dari {total_data_rows} baris untuk efisiensi token."
                        )

                    md_sections.append(
                        f"## Sheet: {ws.title}\n\n" +
                        rows_to_markdown(rows, truncated_note)
                    )
            finally:
                wb.close()

            if not md_sections:
                return ""

            return self._clean_text("\n\n".join(md_sections))

        raise ValueError(f"Format tabular tidak didukung: .{ext}")

    def _looks_like_html(self, text: str) -> bool:
        """Deteksi HTML/webpage dari isi file, bukan hanya ekstensi."""
        if not text:
            return False

        sample = text[:5000].lstrip().lower()
        if sample.startswith('<!doctype html') or sample.startswith('<html'):
            return True

        html_markers = (
            '<head', '<body', '<article', '<main', '<section', '<div',
            '<script', '<style', '<meta', '<link', '<title', '<p', '<h1'
        )
        marker_count = sum(1 for marker in html_markers if marker in sample)
        tag_count = len(re.findall(r'<[a-z][\w:-]*(?:\s[^>]*)?>', sample))
        closing_count = len(re.findall(r'</[a-z][\w:-]*>', sample))

        return marker_count >= 2 or (tag_count >= 8 and closing_count >= 3)

    def _convert_html(self, file_path: str = None, raw_html: str = None) -> str:
        """
        Konversi HTML ke Markdown ringan.
        Bisa menerima file path atau raw HTML string.
        """
        if not BeautifulSoup:
            raise ImportError("Library 'beautifulsoup4' belum terinstall.")

        if raw_html is None:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                raw_html = f.read()

        soup = BeautifulSoup(raw_html, 'html.parser')

        remove_tags = [
            'script', 'style', 'nav', 'header', 'footer', 'aside', 'form',
            'noscript', 'iframe', 'button', 'svg', 'canvas'
        ]
        for tag in soup(remove_tags):
            tag.decompose()

        noisy_keywords = (
            'ad', 'ads', 'advertisement', 'banner', 'billboard', 'gpt',
            'googletag', 'dfp', 'revive', 'share', 'social', 'comment',
            'related', 'recommendation', 'promo', 'popup', 'modal',
            'login', 'subscribe', 'newsletter', 'breadcrumb', 'framebar',
            'navbar', 'footer', 'header'
        )

        for tag in list(soup.find_all(True)):
            if not tag.name:
                continue
            values = []
            tag_id = tag.get('id')
            tag_class = tag.get('class')
            if tag_id:
                values.append(str(tag_id))
            if tag_class:
                values.extend(str(item) for item in tag_class)

            haystack = ' '.join(values).lower()
            if haystack and self._html_attr_has_keyword(haystack, noisy_keywords):
                tag.decompose()

        content_tag = self._pick_html_content_candidate(soup)
        if not content_tag:
            return self._clean_text(soup.get_text(separator='\n'))

        return self._html_tag_to_markdown(content_tag)

    def _pick_html_content_candidate(self, soup):
        """Pilih kandidat artikel dengan teks terlihat terpanjang."""
        candidates = []

        candidates.extend(soup.find_all('article'))
        candidates.extend(soup.find_all('main'))

        content_keywords = ('article', 'detail', 'content', 'read', 'body', 'post')
        for tag in soup.find_all(True):
            values = []
            tag_id = tag.get('id')
            tag_class = tag.get('class')
            if tag_id:
                values.append(str(tag_id))
            if tag_class:
                values.extend(str(item) for item in tag_class)

            haystack = ' '.join(values).lower()
            if haystack and any(keyword in haystack for keyword in content_keywords):
                candidates.append(tag)

        body = soup.find('body')
        if body:
            candidates.append(body)

        if not candidates:
            return None

        unique_candidates = []
        seen = set()
        for candidate in candidates:
            candidate_id = id(candidate)
            if candidate_id not in seen:
                seen.add(candidate_id)
                unique_candidates.append(candidate)

        return max(
            unique_candidates,
            key=lambda tag: len(self._clean_text(tag.get_text(separator='\n')))
        )

    def _html_attr_has_keyword(self, haystack: str, keywords: tuple) -> bool:
        """Cocokkan class/id noise tanpa membuat 'read' match ke keyword pendek 'ad'."""
        tokens = set(filter(None, re.split(r'[^a-z0-9]+', haystack.lower())))
        exact_keywords = {'ad', 'ads', 'gpt', 'dfp'}

        for keyword in keywords:
            if keyword in exact_keywords:
                if keyword in tokens:
                    return True
            elif keyword in haystack:
                return True

        return False

    def _html_tag_to_markdown(self, content_tag) -> str:
        """Konversi subset tag artikel ke teks Markdown hemat token."""
        lines = []

        def append_text(value):
            value = self._clean_text(value)
            if value:
                lines.append(value)

        for tag in content_tag.find_all(
            ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li'],
            recursive=True
        ):
            text = tag.get_text(' ', strip=True)
            if not text:
                continue

            if tag.name.startswith('h') and len(tag.name) == 2:
                level = int(tag.name[1])
                append_text(f"{'#' * level} {text}")
            elif tag.name == 'li':
                append_text(f"* {text}")
            elif tag.name == 'p':
                append_text(text)

        if not lines:
            return self._clean_text(content_tag.get_text(separator='\n'))

        return self._clean_text('\n\n'.join(lines))

    # ----------------------------------------------------------
    # ROUTING & UNIVERSAL FALLBACK
    # ----------------------------------------------------------

    def _universal_fallback(self, file_path: str, ext: str, filename: str) -> str:
        """
        Strategi fallback berlapis untuk file dengan ekstensi tidak dikenal:
        1. Coba baca sebagai teks manusia biasa (UTF-8).
        2. Jika berhasil dan isinya readable â†’ perlakukan sebagai teks.
        3. Jika gagal â†’ bungkus sebagai code block (mungkin source code baru).
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            # Heuristik: jika lebih dari 20% karakter non-printable â†’ biner
            non_printable = sum(1 for c in content if not c.isprintable() and c not in '\n\r\t')
            if len(content) > 0 and (non_printable / len(content)) > 0.20:
                raise UnicodeDecodeError('utf-8', b'', 0, 1, 'Binary content detected')
            if self._looks_like_html(content):
                return self._convert_html(raw_html=content)
            # Konten manusia â†’ perlakukan sebagai teks
            return self._clean_text(content)
        except (UnicodeDecodeError, ValueError):
            # Gagal decode â†’ bungkus sebagai code block
            return self._convert_code(file_path, ext or 'text', filename)

    def process(self, file_path: str, original_filename: str = None) -> dict:
        """
        Fungsi utama â€” dipanggil oleh Frontend/UI.
        Mengembalikan dict: {success, error, data, data_pure, meta}.
        
        Parameter `original_filename` digunakan untuk nama file asli dari user
        (terpisah dari nama temp file di disk).
        """
        # Gunakan nama asli jika disediakan, fallback ke nama file di disk
        display_name = original_filename or os.path.basename(file_path)
        ext = os.path.splitext(display_name)[1].lower().lstrip('.')

        if not os.path.exists(file_path):
            return {
                "success": False,
                "error": "File tidak ditemukan.",
                "data": None,
                "data_pure": None,
                "meta": {
                    "filename": display_name,
                    "estimated_tokens": 0,
                    "estimated_tokens_pure": 0
                }
            }

        try:
            # ---- GUARD: Tolak file gambar ----
            if ext in IMAGE_EXTENSIONS:
                return {
                    "success": False,
                    "error": (
                        f"Format gambar (.{ext}) belum didukung. "
                        "Silakan gunakan file dokumen teks seperti TXT, DOCX, PDF, atau PPTX. "
                        "(OCR direncanakan pada versi berikutnya.)"
                    ),
                    "data": None,
                    "data_pure": None,
                    "meta": {
                        "filename": display_name,
                        "estimated_tokens": 0,
                        "estimated_tokens_pure": 0
                    }
                }

            # ---- ROUTING berdasarkan ekstensi ----
            if ext in ('txt', 'log', 'md'):
                md_result = self._convert_txt(file_path)

            elif ext == 'csv':
                md_result = self._convert_tabular(file_path, 'csv')

            elif ext == 'xlsx':
                md_result = self._convert_tabular(file_path, 'xlsx')

            elif ext == 'docx':
                md_result = self._convert_docx(file_path)

            elif ext == 'pdf':
                md_result = self._convert_pdf(file_path)

            elif ext == 'pptx':
                md_result = self._convert_pptx(file_path)

            elif ext in ('html', 'htm'):
                md_result = self._convert_html(file_path)

            elif ext in CODE_EXTENSIONS:
                md_result = self._convert_code(file_path, ext, display_name)

            else:
                # Universal Fallback untuk ekstensi tidak dikenal
                md_result = self._universal_fallback(file_path, ext, display_name)

            # ---- ULTRA-PURE TEXT MODE: STRIP HEAVY SYMBOLS ----
            # Smart Bypass: abaikan pembersihan jika file adalah source code pemrograman atau web
            is_code = ext in CODE_EXTENSIONS or ext in ('html', 'htm', 'css', 'scss', 'js', 'jsx', 'ts', 'tsx', 'py')
            original_md = md_result
            if not is_code:
                pure_md = self._strip_heavy_symbols(original_md)
            else:
                pure_md = original_md

            return {
                "success": True,
                "error": None,
                "data": original_md,
                "data_pure": pure_md,
                "meta": {
                    "filename": display_name,
                    "estimated_tokens": len(original_md) // 4,
                    "estimated_tokens_pure": len(pure_md) // 4,
                }
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "data": None,
                "data_pure": None,
                "meta": {
                    "filename": display_name,
                    "estimated_tokens": 0,
                    "estimated_tokens_pure": 0
                }
            }

    def _strip_heavy_symbols(self, text: str) -> str:
        """
        Ultra-Pure Text Mode:
        Menghapus simbol hiasan tanpa merusak kata, tabel Markdown,
        atau flowchart ASCII.
        """
        if not text:
            return ""

        # Hapus simbol panah aneh bawaan PowerPoint.
        text = text.replace('\uf0e8', '')
        text = text.replace('\u00ef\u0192\u00a8', '')
        text = text.replace('\x0b', '')
        text = text.replace('\x0c', '')

        # Hapus simbol markdown ringan, tapi JANGAN hapus "-" dan "|".
        text = re.sub(r'[#\*_`]', '', text)

        # Hapus bullet/dekorasi spesifik tanpa global non-ASCII regex.
        text = re.sub(
            r'[\u2022\u2794\u27a2\u27a4\u2713\u2714\u25a0\u27a5\u27a7\u27a8\u27a9\u27aa\u27ab\u27ac\u27ad\u27ae\u27af\u27b1\u27b2\u27bd\u27be]',
            '',
            text
        )

        # Rapikan spasi berlebih.
        text = re.sub(r'[ \t]{2,}', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)

        return text.strip()
